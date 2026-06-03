# -*- coding: utf-8 -*-
"""
================================================================================
PPTX Export Service
================================================================================
Kreira PowerPoint prezentaciju od prevedenih chunk-ova.

Verzija: 1.0.0
================================================================================
"""

import io
import logging
from typing import List, Dict, Any

from pptx import Presentation
from pptx.util import Inches, Pt

from pptx.dml.color import RGBColor

from app.services.base_export_service import BaseExportService

logger = logging.getLogger(__name__)


class PPTXExportService(BaseExportService):
    """Service za export dokumenta u PPTX format."""

    def generate(
        self,
        title: str,
        chunks: List[Dict[str, Any]],
        include_original: bool = False,
    ) -> bytes:
        """
        Generiše PPTX prezentaciju od chunk-ova.

        Args:
            title: Naslov dokumenta
            chunks: Lista chunk-ova (dictionaries)
            include_original: Da li uključuje originalni tekst

        Returns:
            PPTX bytes
        """
        self._report_progress(0, 100, "Priprema PPTX generisanja...")
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        title_slide_layout = prs.slide_layouts[0]
        bullet_slide_layout = prs.slide_layouts[1]

        self._report_progress(10, 100, "Kreiranje naslovnice...")

        title_slide = prs.slides.add_slide(title_slide_layout)
        title_slide.shapes.title.text = title
        subtitle = title_slide.placeholders[1]
        subtitle.text = "AI Learning System - Prevod"

        self._report_progress(20, 100, "Dodavanje slajdova...")
        max_slides = min(len(chunks), 50)

        for i, chunk in enumerate(chunks[:50]):
            translated = self._get_content(chunk)
            original = chunk.get("original_text", "") if include_original else ""

            if len(translated) > 1000:
                translated = translated[:1000] + "..."

            if translated.strip():
                slide = prs.slides.add_slide(bullet_slide_layout)
                shapes = slide.shapes
                title = shapes.title
                body = shapes.placeholders[1]
                title.text = f"Deo {i + 1}"

                tf = body.text_frame
                tf.clear()

                p = tf.paragraphs[0]
                p.text = translated
                p.level = 0
                p.font.size = Pt(18)

                if include_original and original.strip():
                    p = tf.add_paragraph()
                    p.text = f"Original: {original[:500]}..."
                    p.level = 1
                    p.font.size = Pt(12)
                    p.font.color.rgb = RGBColor(128, 128, 128)

            # Progress update
            if i % 10 == 0:
                progress = 20 + int((i / max_slides) * 70)
                self._report_progress(progress, 100, f"Slajd {i + 1}/{max_slides}...")

        self._report_progress(90, 100, "Finalizacija PPTX-a...")
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.font.name = "Arial"

        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)

        self._report_progress(100, 100, "PPTX generisanje završeno!")
        return buffer.getvalue()


pptx_export_service = PPTXExportService()
