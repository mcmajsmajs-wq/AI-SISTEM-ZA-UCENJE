# -*- coding: utf-8 -*-
"""
===============================================================================
ROUNDTRIP EXPORT TESTS (Faza D)
===============================================================================
Verifikuje da export-ovani fajlovi sadrze >= 95% teksta iz originalnih chunkova.

Pokretanje (unutar Docker-a):
    python -m pytest app/tests/unit/test_roundtrip_export.py -v

Svi testovi koriste pytest.importorskip za biblioteke koje nisu dostupne
na host Python-u (fitz, docx, pptx, openpyxl, reportlab).
===============================================================================
"""

import io
import re
import pytest
from typing import List, Dict


# ── Sample test data ──────────────────────────────────────────────────────────

SAMPLE_CHUNKS = [
    {"translated_text": "Prvo Poglavlje", "heading_level": 1},
    {
        "translated_text": (
            "Ovo je prvi pasus u dokumentu. On sluzi kao uvod u temu koja se "
            "obraduje u prvom poglavlju. Sadrzi nekoliko recenica koje opisuju "
            "osnovne koncepte."
        ),
        "heading_level": 0,
    },
    {
        "translated_text": (
            "Drugi pasus nastavlja diskusiju sa detaljnijim objasnjenjem. "
            "Ovde se ulazi u specificne detalje implementacije. Takodje se "
            "pominju vazne napomene i upozorenja."
        ),
        "heading_level": 0,
    },
    {"translated_text": "Drugo Poglavlje — Napredne Teme", "heading_level": 1},
    {"translated_text": "Prva Sekcija", "heading_level": 2},
    {
        "translated_text": (
            "Tekst u prvoj sekciji drugog poglavlja. Objasnjava napredne "
            "koncepte koji nadograduju ono sto je receno u prvom poglavlju. "
            "Preporucuje se da prvo procitate uvod pre nego sto nastavite."
        ),
        "heading_level": 0,
    },
    {"translated_text": "Druga Sekcija", "heading_level": 2},
    {
        "translated_text": (
            "Druga sekcija pokriva dodatne teme. Ovo ukljucuje prakticne "
            "primere i vežbe koje mozete sami da probate. Rezultati se mogu "
            "razlikovati u zavisnosti od vaseg okruzenja."
        ),
        "heading_level": 0,
    },
    {"translated_text": "Podtema 1", "heading_level": 3},
    {
        "translated_text": (
            "Ovo je podtema unutar druge sekcije. Ona ulazi u jos vise "
            "detalja o specificnom aspektu implementacije."
        ),
        "heading_level": 0,
    },
    {"translated_text": "Treca Sekcija", "heading_level": 2},
    {
        "translated_text": (
            "Treca sekcija donosi zakljucak i sumira kljucne tacke. "
            "Ponovo se naglasava vaznost razumevanja osnovnih koncepata "
            "pre nego sto se pređe na napredne teme."
        ),
        "heading_level": 0,
    },
]

SHORT_CHUNKS = [
    {"translated_text": "Kratak naslov", "heading_level": 1},
    {"translated_text": "Kratak tekst.", "heading_level": 0},
]

EMPTY_CHUNKS = []


# ── Helpers ───────────────────────────────────────────────────────────────────


def _input_words(chunks: List[Dict]) -> set:
    """Vraca skup svih reci iz chunkova (heading + body)."""
    words = set()
    for c in chunks:
        text = c.get("translated_text", "") or ""
        words.update(re.findall(r"\w+", text.lower()))
    return words


def _compute_match_percentage(input_chunks: List[Dict], extracted_text: str) -> float:
    """Racuna % reci iz input chunkova koje postoje u ekstraktovanom tekstu."""
    input_w = _input_words(input_chunks)
    if not input_w:
        return 100.0
    output_w = set(re.findall(r"\w+", extracted_text.lower()))
    matched = input_w & output_w
    return len(matched) / len(input_w) * 100


def _extract_pdf_text(pdf_bytes: bytes) -> str:
    """Ekstraktuje sav tekst iz PDF bajtova koristeci PyMuPDF."""
    import fitz

    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    text = ""
    for page in doc:
        text += page.get_text("text") + "\n"
    doc.close()
    return text


def _extract_docx_text(docx_bytes: bytes) -> str:
    """Ekstraktuje sav tekst iz DOCX bajtova."""
    from docx import Document

    doc = Document(io.BytesIO(docx_bytes))
    return "\n".join(p.text for p in doc.paragraphs)


def _extract_pptx_text(pptx_bytes: bytes) -> str:
    """Ekstraktuje sav tekst iz PPTX bajtova."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(pptx_bytes))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
    return "\n".join(texts)


def _extract_xlsx_text(xlsx_bytes: bytes) -> str:
    """Ekstraktuje sav tekst iz XLSX bajtova."""
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(xlsx_bytes))
    texts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        for row in ws.iter_rows(values_only=True):
            texts.extend(str(c) for c in row if c is not None)
    return "\n".join(texts)


# ── PDF Roundtrip Tests ───────────────────────────────────────────────────────


class TestPDFRoundtrip:
    """Roundtrip testovi za PDF export."""

    def test_pdf_roundtrip_full_content_match(self):
        """Svi chunk-ovi sa prevodom se pojavljuju u PDF-u (>=95% reci)."""
        pytest.importorskip("fitz")
        from app.services.pdf_export_service import PDFExportService

        service = PDFExportService()
        pdf_bytes = service.generate(
            title="Roundtrip Test",
            chunks=SAMPLE_CHUNKS,
            target_language="sr",
            include_original=False,
        )

        text = _extract_pdf_text(pdf_bytes)
        match_pct = _compute_match_percentage(SAMPLE_CHUNKS, text)
        assert match_pct >= 95.0, f"Match: {match_pct:.1f}% (expected >= 95%)"

    def test_pdf_roundtrip_short_content(self):
        """Kratki chunk-ovi se pojavljuju u PDF-u."""
        pytest.importorskip("fitz")
        from app.services.pdf_export_service import PDFExportService

        service = PDFExportService()
        pdf_bytes = service.generate(
            title="Short", chunks=SHORT_CHUNKS, target_language="sr"
        )

        text = _extract_pdf_text(pdf_bytes)
        match_pct = _compute_match_percentage(SHORT_CHUNKS, text)
        assert match_pct >= 95.0, f"Short chunks match: {match_pct:.1f}%"

    def test_pdf_roundtrip_include_original(self):
        """include_original=True generise validan PDF (isti content match)."""
        pytest.importorskip("fitz")
        from app.services.pdf_export_service import PDFExportService

        service = PDFExportService()
        pdf_bytes = service.generate(
            title="With Original",
            chunks=SAMPLE_CHUNKS[:2],
            target_language="sr",
            include_original=True,
        )

        text = _extract_pdf_text(pdf_bytes)
        match_pct = _compute_match_percentage(SAMPLE_CHUNKS[:2], text)
        assert match_pct >= 95.0

    def test_pdf_roundtrip_empty_chunks(self):
        """Prazni chunk-ovi rezultiraju validnim PDF-om (0% match = skip)."""
        pytest.importorskip("fitz")
        from app.services.pdf_export_service import PDFExportService

        service = PDFExportService()
        pdf_bytes = service.generate(
            title="Empty", chunks=EMPTY_CHUNKS, target_language="sr"
        )

        assert pdf_bytes.startswith(b"%PDF")
        text = _extract_pdf_text(pdf_bytes)
        assert len(text) > 0  # covers, TOC, footer still present

    def test_pdf_roundtrip_special_chars(self):
        """Specijalni karakteri se pravilno prenose u PDF."""
        pytest.importorskip("fitz")
        from app.services.pdf_export_service import PDFExportService

        special_chunks = [
            {"translated_text": "Č Ć Š Đ Ž — 100% provereno!"},
            {"translated_text": "Cena: $19.99 < 20.00 & povoljno."},
        ]
        service = PDFExportService()
        pdf_bytes = service.generate(
            title="Special", chunks=special_chunks, target_language="sr"
        )

        text = _extract_pdf_text(pdf_bytes)
        match_pct = _compute_match_percentage(special_chunks, text)
        assert match_pct >= 95.0, f"Special chars match: {match_pct:.1f}%"

    def test_pdf_roundtrip_all_headings_preserved(self):
        """Svi heading nivoi se pojavljuju u PDF-u."""
        pytest.importorskip("fitz")
        from app.services.pdf_export_service import PDFExportService

        heading_chunks = [
            {"translated_text": "Naslov 1", "heading_level": 1},
            {"translated_text": "Naslov 2", "heading_level": 2},
            {"translated_text": "Naslov 3", "heading_level": 3},
            {"translated_text": "Ovo je body tekst.", "heading_level": 0},
        ]
        service = PDFExportService()
        pdf_bytes = service.generate(
            title="Headings", chunks=heading_chunks, target_language="sr"
        )

        text = _extract_pdf_text(pdf_bytes)
        for chunk in heading_chunks:
            txt = chunk["translated_text"]
            assert txt in text, f"Heading '{txt}' not found in PDF"

    def test_pdf_roundtrip_dedup_handling(self):
        """Duplikati chunk-ova ne uticu na match score."""
        pytest.importorskip("fitz")
        from app.services.pdf_export_service import PDFExportService

        dup_chunks = SAMPLE_CHUNKS + SAMPLE_CHUNKS  # duplicated
        service = PDFExportService()
        pdf_bytes = service.generate(
            title="Dedup", chunks=dup_chunks, target_language="sr"
        )

        text = _extract_pdf_text(pdf_bytes)
        match_pct = _compute_match_percentage(SAMPLE_CHUNKS, text)
        assert match_pct >= 95.0, f"Dedup match: {match_pct:.1f}%"


# ── DOCX Roundtrip Tests ──────────────────────────────────────────────────────


class TestDOCXRoundtrip:
    """Roundtrip testovi za DOCX export."""

    def test_docx_roundtrip_full_content_match(self):
        """Svi chunk-ovi se pojavljuju u DOCX fajlu (>=95%)."""
        pytest.importorskip("docx")
        from app.services.docx_export_service import DOCXExportService

        service = DOCXExportService()
        docx_bytes = service.generate(
            title="Roundtrip Test",
            chunks=SAMPLE_CHUNKS,
            include_original=False,
        )

        text = _extract_docx_text(docx_bytes)
        match_pct = _compute_match_percentage(SAMPLE_CHUNKS, text)
        assert match_pct >= 95.0, f"Match: {match_pct:.1f}%"

    def test_docx_roundtrip_include_original(self):
        """include_original=True generise validan DOCX (isti content match)."""
        pytest.importorskip("docx")
        from app.services.docx_export_service import DOCXExportService

        service = DOCXExportService()
        docx_bytes = service.generate(
            title="With Original",
            chunks=SAMPLE_CHUNKS[:2],
            include_original=True,
        )

        text = _extract_docx_text(docx_bytes)
        match_pct = _compute_match_percentage(SAMPLE_CHUNKS[:2], text)
        assert match_pct >= 95.0

    def test_docx_roundtrip_empty_chunks(self):
        """Prazni chunk-ovi rezultiraju validnim DOCX-om."""
        pytest.importorskip("docx")
        from app.services.docx_export_service import DOCXExportService

        service = DOCXExportService()
        docx_bytes = service.generate(title="Empty", chunks=EMPTY_CHUNKS)
        assert docx_bytes[:4] == b"PK\x03\x04"

    def test_docx_roundtrip_special_chars(self):
        """Specijalni karakteri se pravilno prenose."""
        pytest.importorskip("docx")
        from app.services.docx_export_service import DOCXExportService

        special_chunks = [
            {"translated_text": "Č Ć Š Đ Ž"},
            {"translated_text": "Cena: $19.99 < 20.00"},
        ]
        service = DOCXExportService()
        docx_bytes = service.generate(title="Special", chunks=special_chunks)

        text = _extract_docx_text(docx_bytes)
        match_pct = _compute_match_percentage(special_chunks, text)
        assert match_pct >= 95.0


# ── PPTX Roundtrip Tests ──────────────────────────────────────────────────────


class TestPPTXRoundtrip:
    """Roundtrip testovi za PPTX export."""

    def test_pptx_roundtrip_full_content_match(self):
        """Svi chunk-ovi se pojavljuju u PPTX fajlu (>=95%)."""
        pytest.importorskip("pptx")
        from app.services.pptx_export_service import PPTXExportService

        service = PPTXExportService()
        pptx_bytes = service.generate(
            title="Roundtrip Test",
            chunks=SAMPLE_CHUNKS,
            include_original=False,
        )

        text = _extract_pptx_text(pptx_bytes)
        match_pct = _compute_match_percentage(SAMPLE_CHUNKS, text)
        assert match_pct >= 95.0

    def test_pptx_roundtrip_slide_limit(self):
        """PPTX limitira na 50 slajdova i dalje sadrzi sav tekst."""
        pytest.importorskip("pptx")
        from app.services.pptx_export_service import PPTXExportService

        many_chunks = [
            {"translated_text": f"Chunk broj {i} sa tekstom."} for i in range(100)
        ]
        service = PPTXExportService()
        pptx_bytes = service.generate(title="Many Chunks", chunks=many_chunks)

        text = _extract_pptx_text(pptx_bytes)
        # Should have text from first 50 chunks (slide limit)
        assert "Chunk broj 0" in text
        assert "Chunk broj 49" in text

    def test_pptx_roundtrip_empty_chunks(self):
        """Prazni chunk-ovi rezultiraju validnim PPTX-om."""
        pytest.importorskip("pptx")
        from app.services.pptx_export_service import PPTXExportService

        service = PPTXExportService()
        pptx_bytes = service.generate(title="Empty", chunks=EMPTY_CHUNKS)
        assert len(pptx_bytes) > 0


# ── XLSX Roundtrip Tests ──────────────────────────────────────────────────────


class TestXLSXRoundtrip:
    """Roundtrip testovi za XLSX export."""

    def test_xlsx_roundtrip_full_content_match(self):
        """Svi chunk-ovi se pojavljuju u XLSX fajlu (>=95%)."""
        pytest.importorskip("openpyxl")
        from app.services.xlsx_export_service import XLSXExportService

        service = XLSXExportService()
        xlsx_bytes = service.generate(
            title="Roundtrip Test",
            chunks=SAMPLE_CHUNKS,
            include_original=False,
        )

        text = _extract_xlsx_text(xlsx_bytes)
        match_pct = _compute_match_percentage(SAMPLE_CHUNKS, text)
        assert match_pct >= 95.0

    def test_xlsx_roundtrip_include_original(self):
        """include_original=True generise validan XLSX (isti content match)."""
        pytest.importorskip("openpyxl")
        from app.services.xlsx_export_service import XLSXExportService

        service = XLSXExportService()
        xlsx_bytes = service.generate(
            title="With Original",
            chunks=SAMPLE_CHUNKS[:2],
            include_original=True,
        )

        text = _extract_xlsx_text(xlsx_bytes)
        match_pct = _compute_match_percentage(SAMPLE_CHUNKS[:2], text)
        assert match_pct >= 95.0

    def test_xlsx_roundtrip_many_chunks(self):
        """Puno chunk-ova se pojavljuje u XLSX-u."""
        pytest.importorskip("openpyxl")
        from app.services.xlsx_export_service import XLSXExportService

        many_chunks = [
            {"translated_text": f"Red {i}: tekst za export."} for i in range(50)
        ]
        service = XLSXExportService()
        xlsx_bytes = service.generate(title="Many Rows", chunks=many_chunks)

        text = _extract_xlsx_text(xlsx_bytes)
        match_pct = _compute_match_percentage(many_chunks, text)
        assert match_pct >= 95.0

    def test_xlsx_roundtrip_empty_chunks(self):
        """Prazni chunk-ovi rezultiraju validnim XLSX-om."""
        pytest.importorskip("openpyxl")
        from app.services.xlsx_export_service import XLSXExportService

        service = XLSXExportService()
        xlsx_bytes = service.generate(title="Empty", chunks=EMPTY_CHUNKS)
        assert xlsx_bytes[:4] == b"PK\x03\x04"
